from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

import io
import re

# ========== THEME DEFINITIONS ==========
THEMES = {
    'professional_blue': {
        'name': 'Professional Blue',
        'bg_start': (227, 242, 253),  # Light blue
        'bg_end': (187, 222, 251),    # Slightly darker blue
        'title_color': (25, 118, 210),  # Deep blue
        'text_color': (33, 33, 33),     # Dark gray
        'accent_color': (66, 165, 245)  # Bright blue
    },
    'modern_dark': {
        'name': 'Modern Dark',
        'bg_start': (48, 48, 48),      # Dark gray
        'bg_end': (33, 33, 33),        # Darker gray
        'title_color': (255, 255, 255), # White
        'text_color': (220, 220, 220),  # Light gray
        'accent_color': (102, 126, 234) # Purple
    },
    'vibrant_orange': {
        'name': 'Vibrant Orange',
        'bg_start': (255, 243, 224),   # Light orange
        'bg_end': (255, 224, 178),     # Darker orange
        'title_color': (230, 81, 0),   # Deep orange
        'text_color': (62, 39, 35),    # Dark brown
        'accent_color': (255, 152, 0)  # Bright orange
    },
    'nature_green': {
        'name': 'Nature Green',
        'bg_start': (232, 245, 233),   # Light green
        'bg_end': (200, 230, 201),     # Darker green
        'title_color': (27, 94, 32),   # Forest green
        'text_color': (33, 33, 33),    # Dark gray
        'accent_color': (76, 175, 80)  # Bright green
    },
    'elegant_purple': {
        'name': 'Elegant Purple',
        'bg_start': (243, 229, 245),   # Light purple
        'bg_end': (225, 190, 231),     # Darker purple
        'title_color': (74, 20, 140),  # Deep purple
        'text_color': (33, 33, 33),    # Dark gray
        'accent_color': (142, 36, 170) # Bright purple
    }
}

# ========== ICON MAPPING ==========
ICON_MAP = {
    # Introduction & Overview
    'introduction': '📊',
    'overview': '👁️',
    'agenda': '📋',
    'outline': '📝',
    
    # Business & Strategy
    'strategy': '🎯',
    'business': '💼',
    'market': '📈',
    'trading': '💹',
    'finance': '💰',
    'investment': '💵',
    'sales': '🤝',
    'marketing': '📢',
    
    # Analysis & Data
    'analysis': '🔍',
    'data': '📊',
    'statistics': '📉',
    'metrics': '📏',
    'report': '📄',
    'research': '🔬',
    
    # Technology
    'technology': '💻',
    'ai': '🤖',
    'artificial intelligence': '🤖',
    'machine learning': '🧠',
    'algorithm': '⚙️',
    'automation': '🔄',
    'digital': '📱',
    'software': '💾',
    
    # Risk & Security
    'risk': '⚠️',
    'security': '🔒',
    'protection': '🛡️',
    'safety': '🦺',
    
    # Growth & Success
    'growth': '📈',
    'success': '🏆',
    'achievement': '🎖️',
    'goals': '🎯',
    'target': '🎯',
    
    # Future & Innovation
    'future': '🔮',
    'innovation': '💡',
    'trends': '📊',
    'forecast': '🌤️',
    'prediction': '🔮',
    
    # Communication
    'communication': '💬',
    'team': '👥',
    'collaboration': '🤝',
    'meeting': '🗓️',
    
    # Process & Timeline
    'process': '⚙️',
    'timeline': '📅',
    'roadmap': '🗺️',
    'workflow': '🔄',
    
    # Results & Conclusion
    'results': '✅',
    'conclusion': '🏁',
    'summary': '📝',
    'takeaway': '🎁',
    'recommendation': '👍',
    
    # Problems & Solutions
    'problem': '❗',
    'challenge': '🧗',
    'solution': '💡',
    'benefits': '✨',
    'advantages': '➕',
    
    # Education & Learning
    'education': '🎓',
    'learning': '📚',
    'training': '🏋️',
    'knowledge': '🧠',
}

# ========== HELPER FUNCTIONS ==========

def get_icon_for_title(title: str) -> str:
    """Find best matching icon for slide title"""
    title_lower = title.lower()
    
    # Check for direct matches first
    for keyword, icon in ICON_MAP.items():
        if keyword in title_lower:
            return icon
    
    # Default icons based on position
    if any(word in title_lower for word in ['intro', 'start', 'begin', 'welcome']):
        return '📊'
    elif any(word in title_lower for word in ['end', 'conclude', 'final', 'summary']):
        return '🏁'
    elif any(word in title_lower for word in ['thank', 'questions', 'q&a']):
        return '🙏'
    
    return ''  # No icon

def clean_text_formatting(text: str) -> str:
    """Remove ** markdown and clean text"""
    # Remove ** for bold
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    
    # Remove * for italic
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    
    # Remove leading bullets if present
    text = re.sub(r'^[\*\-•]\s*', '', text, flags=re.MULTILINE)
    
    return text.strip()

def apply_gradient_background(slide, theme_colors):
    """Apply gradient background to slide"""
    try:
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 45.0  # Diagonal gradient
        
        # Set gradient stops
        fill.gradient_stops[0].color.rgb = RGBColor(*theme_colors['bg_start'])
        fill.gradient_stops[1].color.rgb = RGBColor(*theme_colors['bg_end'])
    except Exception as e:
        print(f"Warning: Could not apply gradient: {e}")
        # Fallback to solid color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*theme_colors['bg_start'])

def add_decorative_bar(slide, theme_colors):
    """Add decorative accent bar at bottom of slide"""
    try:
        left = Inches(0)
        top = slide.shapes[0].top + slide.shapes[0].height - Inches(0.15)
        width = Inches(10)
        height = Inches(0.15)
        
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, width, height
        )
        
        # Style the bar
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*theme_colors['accent_color'])
        shape.line.fill.background()
    except:
        pass  # Skip if error

# ========== MAIN GENERATOR FUNCTION ==========

def generate_pptx(topic: str, sections: list, theme: str = 'professional_blue') -> bytes:
    """Generate a beautifully formatted PowerPoint presentation"""
    
    # Get theme colors
    if theme not in THEMES:
        theme = 'professional_blue'
    theme_colors = THEMES[theme]
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ========== TITLE SLIDE ==========
    title_slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(title_slide_layout)
    
    # Background gradient
    apply_gradient_background(slide, theme_colors)

    # Soft overlay card in center
    try:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.7),
            Inches(1.6),
            Inches(8.6),
            Inches(3.2)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.fill.transparency = 0.15
        card.line.fill.background()
        shadow = card.shadow
        shadow.inherit = False
        shadow.blur_radius = Pt(18)
        shadow.distance = Pt(4)
        shadow.angle = 270.0
        shadow.transparency = 0.6
    except:
        card = None

    # Title text
    left = Inches(1)
    top = Inches(2.1)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = topic
    
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_paragraph.font.size = Pt(50)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(*theme_colors['title_color'])
    
    # Subtitle
    subtitle_top = top + height + Inches(0.2)
    subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, Inches(0.7))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "AI-Generated Presentation"
    
    subtitle_paragraph = subtitle_frame.paragraphs[0]
    subtitle_paragraph.alignment = PP_ALIGN.CENTER
    subtitle_paragraph.font.size = Pt(22)
    subtitle_paragraph.font.color.rgb = RGBColor(*theme_colors['text_color'])
    
    # ========== CONTENT SLIDES ==========
    for i, section in enumerate(sections, 1):
        # Blank layout for full control
        content_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(content_slide_layout)
        
        # Gradient background
        apply_gradient_background(slide, theme_colors)

        # Subtle diagonal overlay
        try:
            overlay = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(-1.0),
                Inches(-0.5),
                Inches(6.0),
                Inches(8.5)
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(*theme_colors['bg_end'])
            overlay.fill.transparency = 0.35
            overlay.rotation = -8
            overlay.line.fill.background()
        except:
            overlay = None
        
        # Rounded title bar with shadow
        try:
            title_bar = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(0.5),
                Inches(0.3),
                Inches(9.0),
                Inches(1.1)
            )
            title_bar.fill.solid()
            title_bar.fill.fore_color.rgb = RGBColor(255, 255, 255)
            title_bar.fill.transparency = 0.05
            title_bar.line.fill.solid()
            title_bar.line.fill.fore_color.rgb = RGBColor(*theme_colors['accent_color'])
            title_bar.line.width = Pt(1.7)

            tb_shadow = title_bar.shadow
            tb_shadow.inherit = False
            tb_shadow.blur_radius = Pt(10)
            tb_shadow.distance = Pt(3)
            tb_shadow.angle = 270.0
            tb_shadow.transparency = 0.6
        except:
            title_bar = None
        
        # Icon for the slide
        icon = get_icon_for_title(section['title'])
        
        # Title text on top of rounded bar
        title_left = Inches(0.7)
        title_top = Inches(0.45)
        title_width = Inches(8.6)
        title_height = Inches(0.9)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_box.text_frame
        
        if icon:
            title_frame.text = f"{icon}  {section['title']}"
        else:
            title_frame.text = section['title']
        
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.font.size = Pt(34)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = RGBColor(*theme_colors['title_color'])
        
        # Thin accent line under header
        try:
            line_shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(0.7),
                Inches(1.4),
                Inches(8.6),
                Inches(0.06)
            )
            line_shape.fill.solid()
            line_shape.fill.fore_color.rgb = RGBColor(*theme_colors['accent_color'])
            line_shape.line.fill.background()
        except:
            pass
        
        # ----- CONTENT AREA -----
        content = section.get('content', '')
        if content:
            # Content card with soft shadow
            try:
                content_bg = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                    Inches(0.8),
                    Inches(1.75),
                    Inches(8.4),
                    Inches(4.8)
                )
                content_bg.fill.solid()
                content_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
                content_bg.fill.transparency = 0.02
                content_bg.line.fill.background()

                c_shadow = content_bg.shadow
                c_shadow.inherit = False
                c_shadow.blur_radius = Pt(14)
                c_shadow.distance = Pt(4)
                c_shadow.angle = 270.0
                c_shadow.transparency = 0.7
            except:
                content_bg = None

            # Text inside content area
            content_left = Inches(1.0)
            content_top = Inches(1.9)
            content_width = Inches(8.0)
            content_height = Inches(4.4)
            
            content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            
            cleaned_content = clean_text_formatting(content)
            lines = [line.strip() for line in cleaned_content.split('\n') if line.strip()]
            
            for j, line in enumerate(lines):
                if j == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                p.text = line
                p.level = 0
                p.font.size = Pt(22)
                p.font.color.rgb = RGBColor(*theme_colors['text_color'])
                p.space_before = Pt(4)
                p.space_after = Pt(4)
                p.line_spacing = 1.2
                p.bullet = True
        
        # Decorative accent bar at bottom (existing helper)
        add_decorative_bar(slide, theme_colors)
        
        # Slide number
        try:
            slide_num_box = slide.shapes.add_textbox(
                Inches(9),
                Inches(7),
                Inches(0.5),
                Inches(0.3)
            )
            slide_num_frame = slide_num_box.text_frame
            slide_num_frame.text = str(i)
            slide_num_para = slide_num_frame.paragraphs[0]
            slide_num_para.font.size = Pt(14)
            slide_num_para.font.color.rgb = RGBColor(*theme_colors['text_color'])
            slide_num_para.alignment = PP_ALIGN.RIGHT
        except:
            pass
    
    # Save to bytes
    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    
    return file_stream.getvalue()
